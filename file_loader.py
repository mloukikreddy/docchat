from pathlib import Path

def load_file(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    if ext == ".txt":
        return _load_txt(file_path)
    elif ext == ".pdf":
        return _load_pdf(file_path)
    elif ext == ".docx":
        return _load_docx(file_path)
    elif ext == ".pptx":
        return _load_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def _load_txt(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def _load_pdf(path):
    from pypdf import PdfReader
    reader = PdfReader(path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)

def _load_docx(path):
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def _load_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text)
    return "\n".join(text)